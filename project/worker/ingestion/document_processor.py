import os
import logging
from typing import Dict, Any
import pdfplumber
import docx
from pptx import Presentation

logger = logging.getLogger(__name__)

class DocumentProcessor:
    def process(self, file_path: str) -> Dict[str, Any]:
        ext = os.path.splitext(file_path)[1].lower()
        if ext == '.pdf':
            return self._process_pdf(file_path)
        elif ext == '.docx':
            return self._process_docx(file_path)
        elif ext == '.pptx':
            return self._process_pptx(file_path)
        return {"error": "Unsupported document type"}

    def _process_pdf(self, path: str) -> Dict[str, Any]:
        text = ""
        meta = {}
        try:
            with pdfplumber.open(path) as pdf:
                meta = pdf.metadata
                for page in pdf.pages:
                    text += page.extract_text() or ""
            return {"type": "pdf", "content": text, "metadata": meta}
        except Exception as e:
            logger.error(f"PDF Processing Error: {e}")
            return {"type": "pdf", "error": str(e)}

    def _process_docx(self, path: str) -> Dict[str, Any]:
        try:
            doc = docx.Document(path)
            text = "\n".join([p.text for p in doc.paragraphs])
            return {"type": "docx", "content": text, "metadata": {}}
        except Exception as e:
            logger.error(f"DOCX Processing Error: {e}")
            return {"type": "docx", "error": str(e)}

    def _process_pptx(self, path: str) -> Dict[str, Any]:
        try:
            prs = Presentation(path)
            text = ""
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + "\n"
            return {"type": "pptx", "content": text, "metadata": {}}
        except Exception as e:
            logger.error(f"PPTX Processing Error: {e}")
            return {"type": "pptx", "error": str(e)}

document_processor = DocumentProcessor()
