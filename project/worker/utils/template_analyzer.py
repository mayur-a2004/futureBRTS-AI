import os
import re
from docx import Document
import json

def analyze_template(doc_path):
    """
    Analyzes an 80-page reference document to extract 'Industrial Gold Standard' formatting and detail density.
    """
    if not os.path.exists(doc_path):
        return {"error": "File not found"}

    try:
        if doc_path.endswith(".docx"):
            doc = Document(doc_path)
            # 1. Extract Headings & Structure
            structure = []
            full_text = []
            for p in doc.paragraphs:
                if p.style.name.startswith('Heading'):
                    structure.append({"level": p.style.name, "text": p.text.strip()})
                if p.text.strip():
                    full_text.append(p.text.strip())
            
            # 2. Calculate Detail Density
            total_words = sum(len(p.split()) for p in full_text)
            avg_density = total_words / len(structure) if structure else 0
            graphics_labels = [p.text for p in doc.paragraphs if "Figure" in p.text or "Table" in p.text]

        elif doc_path.endswith(".pptx"):
            from pptx import Presentation
            prs = Presentation(doc_path)
            structure = []
            full_text = []
            for i, slide in enumerate(prs.slides):
                title = slide.shapes.title.text if slide.shapes.title else f"Slide {i+1}"
                structure.append({"level": "Slide", "text": title})
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        full_text.append(shape.text)
            
            avg_density = sum(len(t.split()) for t in full_text) / len(structure) if structure else 0
            graphics_labels = ["Slide Graphics"] * len(prs.slides)

        elif doc_path.endswith(".zip"):
            import zipfile
            with zipfile.ZipFile(doc_path, 'r') as zip_ref:
                structure = [{"level": "File", "text": n} for n in zip_ref.namelist()[:100]] # Limit to 100 for overview
                full_text = ["ZIP ARCHITECTURE"]
                avg_density = 0
                graphics_labels = ["Source Code Files"]
        else:
            return {"error": "Unsupported format"}

        blueprint = {
            "template_name": os.path.basename(doc_path),
            "type": "PPTX" if doc_path.endswith(".pptx") else "DOCX",
            "detected_points": len(structure),
            "structure": structure[:50],
            "avg_words_per_section": int(avg_density),
            "required_visuals": list(set([re.sub(r'[^a-zA-Z\s]', '', l).strip() for l in graphics_labels[:10]])),
            "status": "VALIDATED_MASTER_BLUEPRINT"
        }

        # Save to storage for Pipeline to use
        output_path = os.path.join("storage", "templates", "master_blueprint.json")
        os.makedirs(os.path.dirname(output_path), exist_ok=True)
        with open(output_path, "w") as f:
            json.dump(blueprint, f, indent=4)

        return blueprint

    except Exception as e:
        return {"error": str(e)}

if __name__ == "__main__":
    import sys
    if len(sys.argv) > 1:
        print(json.dumps(analyze_template(sys.argv[1]), indent=2))
