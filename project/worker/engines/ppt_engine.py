from pptx import Presentation
import os

def generate_ppt(payload):
    os.makedirs("outputs", exist_ok=True)
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[1]) # Title and Content
    slide.shapes.title.text = "Future BRTS"
    slide.placeholders[1].text = payload.prompt or ""

    path = f"outputs/ppt_{payload.taskId}.pptx"
    prs.save(path)

    return {"status": "SUCCESS", "artifacts": {"ppt": path}}
