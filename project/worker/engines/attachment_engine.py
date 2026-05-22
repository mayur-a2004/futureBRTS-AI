import os
import pdfplumber
from pptx import Presentation
from docx import Document
# from PIL import Image # Optional for basic image info

def process_attachment(payload):
    """
    Process an attachment based on its type.
    Payload (WorkerRequest) -> metadata contains:
      - filePath: Absolute path to the file
      - fileType: 'application/pdf', 'image/png', etc.
      - originalName: Original filename
    """
    try:
        data = payload.metadata
        file_path = data.get('filePath')
        file_type = data.get('fileType', '').lower()
        
        if not file_path or not os.path.exists(file_path):
            return {"status": "FAIL", "reason": f"File not found at path: {file_path}"}

        extracted_text = ""
        summary = ""
        
        # 1. PDF
        if "pdf" in file_type:
            with pdfplumber.open(file_path) as pdf:
                for page in pdf.pages:
                    extracted_text += (page.extract_text() or "") + "\n"
            summary = f"Extracted {len(extracted_text)} characters from PDF."

        # 2. DOCX
        elif "word" in file_type or "docx" in file_type:
            doc = Document(file_path)
            for para in doc.paragraphs:
                extracted_text += para.text + "\n"
            summary = f"Extracted {len(extracted_text)} characters from Word Document."
            
        # 3. PPTX
        elif "presentation" in file_type or "pptx" in file_type:
            prs = Presentation(file_path)
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        extracted_text += shape.text + "\n"
            summary = f"Extracted {len(extracted_text)} characters from Presentation."
            
        # 4. Text / Code
        elif "text" in file_type or file_path.endswith((".txt", ".py", ".js", ".ts", ".md", ".json")):
            with open(file_path, 'r', encoding='utf-8') as f:
                extracted_text = f.read()
            summary = f"Extracted {len(extracted_text)} characters from Text file."

        # 5. Image (Stub)
        elif "image" in file_type:
            # Placeholder for OCR
            summary = "Image file received. Visual content is ready for review."
            extracted_text = f"USER_UPLOADED_IMAGE: {data.get('originalName', 'Unknown')}. (Note: Deep visual analysis is pending integration, treat this as a visual reference provided by the user)."

        else:
            summary = f"Unsupported file type: {file_type}"
            extracted_text = "[UNSUPPORTED TYPE]"

        return {
            "status": "SUCCESS",
            "summary": summary,
            "extracted_text": extracted_text,
            "file_path": file_path
        }

    except Exception as e:
        print(f"Attachment Processing Error: {e}")
        return {"status": "FAIL", "reason": str(e)}
