import os
import json
import re
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from utils.diagram_engine import diagram_engine
from utils.sanitizer import sanitize_for_xml

PROJECTS_ROOT = os.path.abspath(os.path.join(os.path.dirname(__file__), "..", "..", "backend", "projects"))
os.makedirs(PROJECTS_ROOT, exist_ok=True)

def sanitize_for_xml(text):
    """Removes control characters that break XML-based formats (DOCX/PPTX)."""
    if not text: return ""
    illegal_chars = re.compile(r'[\x00-\x08\x0b\x0c\x0e-\x1f\x7f-\x84\x86-\x9f\ud800-\udfff\ufdd0-\ufdef\ufffe\uffff]')
    return illegal_chars.sub('', str(text))

def generate_project_ppt(project_name, category, field, project_type, tech_stack, vision, ai_content, job_id, images=None):
    prs = Presentation()

    def add_slide(title, content, image_path=None):
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        title_shape = slide.shapes.title
        title_shape.text = title
        
        body_shape = slide.shapes.placeholders[1]
        tf = body_shape.text_frame
        tf.word_wrap = True
        
        # Split content into bullets
        if "\n" in content:
            for line in content.split("\n"):
                if line.strip():
                    p = tf.add_paragraph()
                    p.text = line.strip()
                    p.level = 0
        else:
            tf.text = content

        if image_path:
            try:
                # Add image on the right or bottom
                slide.shapes.add_picture(image_path, Inches(5.5), Inches(2), width=Inches(4))
            except: pass

    # Load Master Blueprint (Academic Standard) - Fix 5: Absolute Path
    base_dir = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
    blueprint_path = os.path.join(base_dir, "storage", "templates", "master_blueprint.json")
    master_blueprint = {}
    if os.path.exists(blueprint_path):
        import json
        with open(blueprint_path, "r") as f:
            master_blueprint = json.load(f)
    else:
        print(f"⚠️ MASTER BLUEPRINT NOT FOUND AT {blueprint_path}")

    # 2. INDUSTRIAL BOILERPLATE ENGINE
    academic_boilerplate = {
        "Abstract": f"Project Mission: {vision}\n- Tactical 3-tier architecture implementation.\n- High-fidelity industrial synchronization.\n- Production-ready {tech_stack} deployment.",
        "System Requirements": "HARDWARE:\n- Minimum 8GB RAM / Quad-core CPU.\n- SSD Optimized Architecture.\nSOFTWARE:\n- Node.js Runtime Environment.\n- Secure Database persistence.",
        "System Design": "ARCHITECTURE OVERVIEW:\n- Decoupled Frontend/Backend nodes.\n- RESTful API Orchestration.\n- 3-Layer Security Handshake.",
        "Database Design": "SCHEMA STRATEGY:\n- 3NF Normalized Models.\n- ACID Compliance Verified.\n- Automated Data Backups.",
        "Future Scope": "ROADMAP:\n- AI/ML Predictive Analytics.\n- Cross-Platform UI Expansion.\n- Microservices Scalability."
    }

    # helper for AI extraction (More robust)
    def get_section_text(num, title, fallback):
        pattern = rf"(?:{num}\.?\s+|[Cc]hapter\s*{num}\.?\s+){title}(.*?)(?=(?:\d+\.?\s+|[Cc]hapter\s*\d+\.?\s+|$))"
        match = re.search(pattern, ai_content, re.DOTALL | re.IGNORECASE)
        
        extracted_content = fallback
        mermaid_codes = []

        if match:
             extracted_content = match.group(1).strip().replace("*", "").replace("- ", "\n- ").strip()
        
        # SEARCH FOR MERMAID BLOCKS
        mermaid_pattern = r"\[MERMAID_START\](.*?)\[MERMAID_END\]"
        matches = re.findall(mermaid_pattern, ai_content if not match else match.group(1), re.DOTALL)
        for m in matches:
            mermaid_codes.append(m.strip())

        return extracted_content, mermaid_codes

    # 1. Title Slide (Industrial Layout)
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    title = slide.shapes.title
    title.text = project_name.upper()
    
    if images and 'logo' in images:
        try:
            slide.shapes.add_picture(images['logo'], Inches(0.5), Inches(0.5), width=Inches(1.5))
        except: pass

    subtitle = slide.shapes.placeholders[1]
    subtitle.text = f"TITAN SUPREME MISSION: {field}\nIndustrial Build for {category}\nArchitected by Genesis Neural Core"

    # Slide Selection Logic
    slides_to_render = master_blueprint.get("structure", [])
    if not slides_to_render:
        slides_to_render = [
            {"text": "Abstract", "num": "4"},
            {"text": "Introduction", "num": "5"},
            {"text": "System Requirements", "num": "9"},
            {"text": "System Design", "num": "12"},
            {"text": "Database Design", "num": "13"},
            {"text": "Module Description", "num": "15"},
            {"text": "Visual Representation", "num": "19"},
            {"text": "Future Scope", "num": "21"}
        ]

    for i, section in enumerate(slides_to_render[:10], 1): # Limit to 10 for presentation quality
        title = section.get("text", "Module")
        num = section.get("num", str(i))
        fallback = academic_boilerplate.get(title, f"In-depth technical analysis and industrial roadmap for {project_name} following the {tech_stack} architecture.")
        
        content_text, mermaid_diagrams = get_section_text(num, title, fallback)
        
        # Slide with text
        img = None
        if ("Visual" in title or "Design" in title or num == "19") and images and 'mockup' in images:
            img = images['mockup']
            
        add_slide(sanitize_for_xml(title), sanitize_for_xml(content_text if len(content_text) > 10 else fallback), img)

        # Slide(s) with diagrams
        diagram_dir = os.path.join("storage", "projects", job_id, "diagrams")
        os.makedirs(diagram_dir, exist_ok=True)
        for idx, m_code in enumerate(mermaid_diagrams):
             try:
                  img_path = diagram_engine.render_mermaid(m_code, output_dir=diagram_dir)
                  if img_path:
                       add_slide(f"{sanitize_for_xml(title)} - Visual Model {idx+1}", "Industrial Schematic Representation", img_path)
             except Exception as e:
                  print(f"PPT Diagram Render Failed: {e}")

    # Final Slide
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    title_final = slide.shapes.title
    title_final.text = "Thank You"
    subtitle = slide.shapes.placeholders[1]
    subtitle.text = "TITAN BUILD COMPLETE\nQuestions & Feedback"

    safe_name = re.sub(r'[^\w\s-]', '', project_name).strip().replace(' ', '_')
    filename = f"Presentation_{safe_name}_{job_id[:8]}.pptx"
    ppt_dir = os.path.join(PROJECTS_ROOT, job_id, "docs")
    os.makedirs(ppt_dir, exist_ok=True)
    path = os.path.join(ppt_dir, filename)
    prs.save(path)
    return path
