import re
import os
from typing import Dict, Any
from rich.console import Console

console = Console()

class TextBooster:
    @staticmethod
    def analyze(text: str) -> Dict[str, Any]:
        """
        Instant NLP Analysis (Standard Lib Version).
        Reliable & Fast.
        """
        if not text: return {}
        
        words = text.split()
        word_count = len(words)
        sentences = re.split(r'[.!?]+', text)
        sentence_count = len([s for s in sentences if s.strip()])
        
        # Simple Logic for "Boost" appearance without heavy deps
        avg_word_len = sum(len(w) for w in words) / word_count if word_count > 0 else 0
        
        return {
            "analysis_type": "Statistical (Fast)",
            "word_count": word_count,
            "sentence_count": sentence_count,
            "avg_word_length": round(avg_word_len, 2),
            "language_detected": "en" # Placeholder for stability
        }

def run(doc_path: str, command_info: Dict[str, Any]) -> Dict[str, Any]:
    """
    Supercharged Document Pipeline.
    Supports txt, docx (try/except), pdf (try/except).
    """
    extracted_text = ""
    meta = {}
    
    # 1. READ FAST with Fallbacks
    try:
        if doc_path.endswith(".pdf"):
            try:
                import pdfplumber
                with pdfplumber.open(doc_path) as pdf:
                    pages = []
                    for p in pdf.pages:
                        extracted = p.extract_text()
                        if extracted: pages.append(extracted)
                    extracted_text = "\n".join(pages)
                    meta["pages"] = len(pdf.pages)
            except ImportError:
                 extracted_text = "[Error: pdfplumber not installed]"

        elif doc_path.endswith(".docx"):
            try:
                import docx
                doc = docx.Document(doc_path)
                extracted_text = "\n".join([p.text for p in doc.paragraphs])
            except ImportError:
                extracted_text = "[Error: python-docx not installed]"
            except Exception as e:
                extracted_text = f"[Error reading docx: {e}]"
        elif doc_path.endswith(".xlsx") or doc_path.endswith(".xls"):
            try:
                import pandas as pd
                df = pd.read_excel(doc_path)
                # Convert first few rows to markdown for easy reading by LLM
                extracted_text = df.head(50).to_markdown(index=False)
                meta["rows"] = len(df)
                meta["columns"] = list(df.columns)
            except ImportError:
                 extracted_text = "[Error: pandas/openpyxl not installed]"
            except Exception as e:
                 extracted_text = f"[Error reading excel: {e}]"

        elif doc_path.endswith(".csv"):
             try:
                import pandas as pd
                # from ml.data_scientist import ds_engine
                
                df = pd.read_csv(doc_path)
                
                # Basic Extract for LLM context
                extracted_text = df.head(50).to_markdown(index=False)
                # meta.update(ds_engine.analyze_dataset(df)) # Optional deep analysis
                meta["rows"] = len(df)

             except Exception as e:
                 extracted_text = f"[Error reading csv: {e}]"
                 
        else:
            # Fallback text
            try:
                with open(doc_path, "r", encoding="utf-8") as f:
                    extracted_text = f.read()
            except Exception as e:
                extracted_text = f"[Error reading text file: {e}]"
    except Exception as e:
        print(f"CRITICAL ERROR READ FILE: {e}")
        extracted_text = f"Error reading file: {str(e)}"

    # 2. ANALYZE FAST
    if not extracted_text:
        extracted_text = "No text extracted."
        
    nlp_result = TextBooster.analyze(extracted_text[:5000]) 
    
    return {
        "status": "completed",
        "extracted_text": extracted_text,
        "result": {
            "nlp_analysis": nlp_result,
            "metadata": meta
        }
    }

def create_docx(content: str, params: dict = {}) -> dict:
    console.print(f"[bold blue]📄 DOCX CREATION TRIGGERED[/bold blue]")
    try:
        from docx import Document
    except ImportError:
        return {"status": "failed", "error": "python-docx not installed"}

    output_dir = "storage/projects/generated_docs"
    os.makedirs(output_dir, exist_ok=True)
    filename = f"doc_{os.urandom(4).hex()}.docx"
    output_path = os.path.join(output_dir, filename)
    download_url = f"/downloads/generated_docs/{filename}"

    try:
        document = Document()
        document.add_heading("Generated Document", 0)
        
        # Split content by newlines and add paragraphs
        for line in content.split('\n'):
            if line.strip():
                document.add_paragraph(line)
                
        document.save(output_path)
        
        return {
            "status": "completed",
            "extracted_text": f"Docx created with content length: {len(content)}",
            "result": {
                "file_path": output_path,
                "download_url": download_url,
                "content_preview": content[:100] + "..."
            }
        }
    except Exception as e:
        return {"status": "failed", "error": str(e)}

def create_ppt(slides: list, params: dict = {}) -> dict:
    console.print(f"[bold orange1]📊 PPT CREATION TRIGGERED[/bold orange1]")
    try:
        from pptx import Presentation
        from pptx.util import Inches
    except ImportError:
         return {"status": "failed", "error": "python-pptx not installed"}

    output_dir = "storage/projects/generated_docs"
    os.makedirs(output_dir, exist_ok=True)
    filename = f"presentation_{os.urandom(4).hex()}.pptx"
    output_path = os.path.join(output_dir, filename)
    download_url = f"/downloads/generated_docs/{filename}"

    try:
        prs = Presentation()
        
        for slide_data in slides:
            # Add slide (Layout 1 = Title + Content)
            slide = prs.slides.add_slide(prs.slide_layouts[1])
            title = slide.shapes.title
            content = slide.placeholders[1]
            
            # Allow string slides or object slides
            if isinstance(slide_data, str):
                title.text = "Slide"
                content.text = slide_data
            else:
                title.text = slide_data.get("title", "Slide")
                content.text = slide_data.get("content", "")

        prs.save(output_path)
        
        return {
            "status": "completed",
            "extracted_text": f"Presentation created with {len(slides)} slides.",
            "result": {
                "file_path": output_path,
                "download_url": download_url
            }
        }
    except Exception as e:
        return {"status": "failed", "error": str(e)}

def convert(file_path: str, to_format: str) -> dict:
    console.print(f"[bold green]🔄 CONVERT TRIGGERED: {to_format}[/bold green]")
    # Placeholder for simple conversion
    # Only supporting Text -> Docx for now as robust example
    
    if to_format == "docx":
        # Extract text using run() then create_docx logic
        analysis = run(file_path, {})
        text = analysis.get("extracted_text", "")
        return create_docx(text)
    
    return {"status": "failed", "error": "Conversion format not supported"}
